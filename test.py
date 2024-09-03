from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_field_names(pptx_file):
    prs = Presentation(pptx_file)
    fields = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.has_text_frame:
                text = shape.text.strip()
                if text.startswith('{') and text.endswith('}'):
                    fields.append(text[1:-1])

    return fields

# Пример использования
pptx_file = 'files/Сертификат_поля_печать.pptx'
field_names = extract_field_names(pptx_file)
print("Найденные имена полей:")
for name in field_names:
    print(name)