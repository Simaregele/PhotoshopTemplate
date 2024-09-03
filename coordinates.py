from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text_coordinates(pptx_file):
    presentation = Presentation(pptx_file)
    text_objects = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or \
               (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.text):
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                text = shape.text
                text_objects.append({
                    'text': text,
                    'coordinates': (left, top, width, height)
                })

    return text_objects

# Пример использования
pptx_file = 'files/Сертификат_поля_печать.pptx'
result = extract_text_coordinates(pptx_file)

for obj in result:
    print(f"Текст: {obj['text']}")
    print(f"Координаты: {obj['coordinates']}")
    print("---")