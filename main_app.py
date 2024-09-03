import streamlit as st
import json
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

# Путь к шаблону
TEMPLATE_PATH = os.path.join('files', 'Сертификат_поля_печать.pptx')


def load_coordinates(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def find_shape_by_coordinates(slide, target_x, target_y, tolerance=100000):
    for shape in slide.shapes:
        if abs(shape.left - target_x) < tolerance and abs(shape.top - target_y) < tolerance:
            return shape
    return None


def create_new_pptx(fields_data, coordinates, output_path):
    prs = Presentation(TEMPLATE_PATH)
    for slide in prs.slides:
        for field, coords in coordinates.items():
            shape = find_shape_by_coordinates(slide, coords['x'], coords['y'])
            if shape and shape.has_text_frame:
                text_frame = shape.text_frame
                if text_frame.text:
                    text_frame.text = fields_data[field]
                else:
                    p = text_frame.paragraphs[0]
                    p.text = fields_data[field]

    prs.save(output_path)


def main():
    st.title("PowerPoint Field Editor")

    # Проверка наличия файла шаблона
    if not os.path.exists(TEMPLATE_PATH):
        st.error(f"Шаблон PowerPoint не найден по пути: {TEMPLATE_PATH}")
        return

    # Load coordinates
    coordinates = load_coordinates('field_coordinates.json')

    # Create a text input for each field
    fields_data = {}
    for field in coordinates.keys():
        fields_data[field] = st.text_input(f"{field.replace('_', ' ').title()}")

    # Button to create new PowerPoint
    if st.button("Create New PowerPoint"):
        # Create new PowerPoint file
        output_path = "new_powerpoint.pptx"
        create_new_pptx(fields_data, coordinates, output_path)

        # Provide download link for new file
        with open(output_path, "rb") as file:
            btn = st.download_button(
                label="Download New PowerPoint",
                data=file,
                file_name="new_powerpoint.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                on_click=lambda: os.remove(output_path)
            )

        st.success("New PowerPoint created successfully!")


if __name__ == "__main__":
    main()